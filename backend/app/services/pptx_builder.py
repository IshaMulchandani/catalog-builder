import os
from io import BytesIO
from typing import Optional

import httpx
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

from ..schemas import SlidePlan
from .image_service import UPLOAD_DIR, is_remote_url

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


# How much darker the "Model" subtitle color is than the accent color it's
# derived from — shared constant with the web preview's darkenHex() so both
# renderers land on the same shade. 0.25 = each channel scaled to 75% of its
# original value (toward black).
MODEL_DARKEN_AMOUNT = 0.25


def _darken_rgb(rgb: RGBColor, amount: float = MODEL_DARKEN_AMOUNT) -> RGBColor:
    r, g, b = rgb[0], rgb[1], rgb[2]
    scale = 1 - amount
    return RGBColor(int(r * scale), int(g * scale), int(b * scale))


EMU_PER_INCH = 914400


def _estimate_line_count(text: str, box_width_emu: int, font_pt: float, max_lines: int = 2) -> int:
    """Rough estimate of how many lines `text` needs to wrap to inside a box
    of the given width at the given font size. Used to size the "Model"
    subtitle's reserved height to what it actually needs — a fixed
    always-2-lines reservation left a visible dead gap above the divider
    for short values (PowerPoint text boxes, unlike CSS, can't auto-hug
    their own content height while also auto-shrinking on overflow, so this
    has to be estimated up front rather than measured)."""
    if not text:
        return 1
    box_width_in = box_width_emu / EMU_PER_INCH
    avg_char_width_in = (font_pt * 0.52) / 72  # rough glyph width for a semibold sans font
    chars_per_line = max(1, int(box_width_in / avg_char_width_in))
    lines = -(-len(text) // chars_per_line)  # ceil division
    return max(1, min(max_lines, lines))


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def _load_image_bytes(image_path: str) -> Optional[bytes]:
    """Load raw image bytes given either a full URL (R2, production) or a
    local filename under UPLOAD_DIR (local dev). Returns None if the image
    can't be fetched/found, so callers can skip it gracefully."""
    if is_remote_url(image_path):
        try:
            resp = httpx.get(image_path, timeout=10.0, follow_redirects=True)
            resp.raise_for_status()
            return resp.content
        except Exception:
            return None
    local_path = os.path.join(UPLOAD_DIR, image_path)
    if os.path.exists(local_path):
        with open(local_path, "rb") as f:
            return f.read()
    return None


def _add_cover_image(slide, img_bytes: bytes, x, y, w, h):
    """Place an image filling exactly the (x, y, w, h) box, cropping the source
    image (never stretching it) to emulate CSS `object-fit: cover` — the same
    behavior the web preview uses. Falls back to a plain (stretched) placement
    if the image can't be read."""
    pic = slide.shapes.add_picture(BytesIO(img_bytes), x, y, width=w, height=h)
    try:
        with Image.open(BytesIO(img_bytes)) as im:
            src_w, src_h = im.size
    except Exception:
        return pic

    if not src_w or not src_h:
        return pic

    img_ratio = src_w / src_h
    box_ratio = w / h

    if img_ratio > box_ratio:
        # source is relatively wider than the target box: crop the left/right edges
        crop_frac = 1 - (box_ratio / img_ratio)
        pic.crop_left = crop_frac / 2
        pic.crop_right = crop_frac / 2
    elif img_ratio < box_ratio:
        # source is relatively taller than the target box: crop the top/bottom edges
        crop_frac = 1 - (img_ratio / box_ratio)
        pic.crop_top = crop_frac / 2
        pic.crop_bottom = crop_frac / 2

    return pic


def _add_contain_image(slide, img_bytes: bytes, x, y, w, h):
    """Place an image within the (x, y, w, h) box, preserving its aspect ratio
    and centering it (never stretching) — matches CSS `object-fit: contain`,
    which is what the logo overlay in the web preview uses. Distortion is
    fine for cropped product photos but not for a logo."""
    try:
        with Image.open(BytesIO(img_bytes)) as im:
            src_w, src_h = im.size
    except Exception:
        return slide.shapes.add_picture(BytesIO(img_bytes), x, y, width=w, height=h)

    if not src_w or not src_h:
        return slide.shapes.add_picture(BytesIO(img_bytes), x, y, width=w, height=h)

    img_ratio = src_w / src_h
    box_ratio = w / h

    if img_ratio > box_ratio:
        draw_w = w
        draw_h = int(w / img_ratio)
    else:
        draw_h = h
        draw_w = int(h * img_ratio)

    draw_x = x + (w - draw_w) // 2
    draw_y = y + (h - draw_h) // 2
    return slide.shapes.add_picture(BytesIO(img_bytes), draw_x, draw_y, width=draw_w, height=draw_h)


def _add_new_badge(slide, card_x, card_y, card_w, accent_rgb: RGBColor):
    """Small pill badge pinned to the top-right corner of a product card,
    marking products created within the last NEW_BADGE_DAYS days."""
    badge_w = Inches(0.62)
    badge_h = Inches(0.28)
    inset = Inches(0.08)
    bx = card_x + card_w - badge_w - inset
    by = card_y + inset

    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, badge_w, badge_h)
    badge.adjustments[0] = 0.5  # fully rounded ends (pill shape)
    badge.fill.solid()
    badge.fill.fore_color.rgb = accent_rgb
    badge.line.fill.background()
    badge.shadow.inherit = False

    tf = badge.text_frame
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "NEW"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return badge


def _add_accent_bar(slide, accent_rgb: RGBColor):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.12), SLIDE_H)  # rectangle
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_rgb
    bar.line.fill.background()


def _add_text(slide, left, top, width, height, text, size, bold=False, color=None, align=PP_ALIGN.LEFT, auto_fit=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if auto_fit:
        # Tells PowerPoint (and LibreOffice, used for PDF export) to shrink
        # the font to fit this box when the file is opened/rendered — the
        # same "shrink text on overflow" behavior the web preview does live,
        # so long product descriptions never get truncated in the export
        # either, no matter how long they are.
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return box


def build_pptx(plan: SlidePlan, accent_color: str, currency_symbol: str) -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    accent_rgb = _hex_to_rgb(accent_color)

    for slide_desc in plan.slides:
        slide = _blank_slide(prs)
        _add_accent_bar(slide, accent_rgb)

        if slide_desc.type == "cover":
            _add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(1.2),
                       slide_desc.title or "", 44, bold=True)
            _add_text(slide, Inches(1), Inches(3.9), Inches(11), Inches(0.8),
                       slide_desc.subtitle or "", 20, color=RGBColor(0x66, 0x66, 0x66))

            if slide_desc.logo_path:
                logo_bytes = _load_image_bytes(slide_desc.logo_path)
                if logo_bytes:
                    lx = slide_desc.logo_x if slide_desc.logo_x is not None else 0.75
                    ly = slide_desc.logo_y if slide_desc.logo_y is not None else 0.06
                    lw = slide_desc.logo_w if slide_desc.logo_w is not None else 0.18
                    lh = slide_desc.logo_h if slide_desc.logo_h is not None else 0.18
                    _add_contain_image(
                        slide, logo_bytes,
                        int(SLIDE_W * lx), int(SLIDE_H * ly),
                        int(SLIDE_W * lw), int(SLIDE_H * lh),
                    )

        elif slide_desc.type == "index":
            _add_text(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
                       "Index", 32, bold=True)

            # Categories used to be stacked in a single column that ran off
            # the bottom of the slide once there were more than ~7-8 of
            # them (silently lost -- nothing rendered past the slide edge).
            # Now they flow into as many columns as needed, each separated
            # by a thin vertical rule, filling one column top-to-bottom
            # before starting the next.
            categories = slide_desc.categories or []
            content_left = Inches(1)
            content_width = Inches(11)
            y_start = Inches(1.8)
            row_h = Inches(0.7)
            bottom_margin = Inches(0.5)
            col_gap = Inches(0.4)

            # Column count is driven purely by how many rows fit vertically --
            # no cap on column count. However many columns that implies, they
            # split the fixed content_width evenly (see col_width below), so
            # more categories means narrower columns rather than anything
            # overflowing past the slide edge.
            usable_height = SLIDE_H - bottom_margin - y_start
            rows_per_col = max(1, int(usable_height // row_h))
            num_cols = max(1, -(-len(categories) // rows_per_col)) if categories else 1

            col_width = (
                (content_width - col_gap * (num_cols - 1)) // num_cols
                if num_cols > 1 else content_width
            )

            for col in range(num_cols):
                col_x = content_left + col * (col_width + col_gap)
                chunk = categories[col * rows_per_col: (col + 1) * rows_per_col]
                y = y_start
                for name in chunk:
                    _add_text(slide, col_x, y, col_width, row_h, name, 22, auto_fit=True)
                    y += row_h

                if col < num_cols - 1:
                    divider_x = col_x + col_width + col_gap // 2
                    divider = slide.shapes.add_shape(1, divider_x, y_start, Pt(1), usable_height)
                    divider.fill.solid()
                    divider.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
                    divider.line.fill.background()
                    divider.shadow.inherit = False

        elif slide_desc.type == "category":
            title = slide_desc.title or ""
            # Raised closer to the top edge (was y=0.4) to open up breathing room
            # between the title and the product cards below it.
            title_top = Inches(0.22)
            title_h = Inches(0.7)
            _add_text(slide, Inches(0.6), title_top, Inches(11), title_h, title, 28, bold=True)

            # 2x2 grid of horizontal cards: image on the left, details on the right.
            # Vertically centered in the space below the title (equal gap above
            # and below the grid) instead of pinned at a fixed offset -- the old
            # fixed y-positions left a big gap under the title and no gap at all
            # above the slide's bottom edge.
            products = slide_desc.products or []
            cell_w = Inches(5.9)
            cell_h = Inches(2.9)
            row_gap = Inches(0.2)
            grid_h = cell_h * 2 + row_gap

            title_bottom = title_top + title_h
            leftover = (SLIDE_H - title_bottom) - grid_h
            vertical_gap = max(leftover // 2, 0)
            grid_top = title_bottom + vertical_gap

            xs = [Inches(0.7), Inches(6.9)]
            ys = [grid_top, grid_top + cell_h + row_gap]
            positions = [(xs[0], ys[0]), (xs[1], ys[0]), (xs[0], ys[1]), (xs[1], ys[1])]

            img_w = Inches(2.4)  # ~42% of cell_w, matching the web preview's image column
            text_left_offset = Inches(2.55)
            text_w = cell_w - text_left_offset
            border_inset = Inches(0.03)  # keeps the card border visible around the image edge

            for product, (x, y) in zip(products, positions):
                card = slide.shapes.add_shape(1, x, y, cell_w, cell_h)  # rectangle outline
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                card.line.color.rgb = RGBColor(0xE5, 0xE5, 0xE5)
                card.line.width = Pt(0.75)
                card.shadow.inherit = False

                if product.image_path:
                    img_bytes = _load_image_bytes(product.image_path)
                    if img_bytes:
                        _add_cover_image(
                            slide, img_bytes,
                            x + border_inset, y + border_inset,
                            img_w - border_inset, cell_h - 2 * border_inset,
                        )

                # Category name is already the slide title, so it isn't repeated
                # per-card — the brand/product name is now the card's headline,
                # in the accent color, taking the vertical space the category
                # eyebrow used to occupy.
                text_x = x + text_left_offset
                name_top = y + Inches(0.18)
                name_h = Inches(0.4)
                _add_text(slide, text_x, name_top, text_w, name_h,
                           product.name, 18, bold=True, color=accent_rgb)

                # Optional "Model" subtitle — smaller than the brand name,
                # larger than the description, in a darker shade of the
                # accent color. Only takes up space (and only shifts the
                # divider/description down) when actually present, so
                # products without one render exactly as before.
                model_text = (product.model or "").strip()
                divider_top = name_top + name_h + Inches(0.10)
                if model_text:
                    model_top = name_top + name_h
                    model_font_pt = 14
                    model_lines = _estimate_line_count(model_text, text_w, model_font_pt)
                    model_h = Inches(0.24) * model_lines  # sized to the estimated line count, not always 2
                    _add_text(slide, text_x, model_top, text_w, model_h,
                               model_text, model_font_pt, color=_darken_rgb(accent_rgb), auto_fit=True)
                    divider_top = model_top + model_h + Inches(0.06)

                divider = slide.shapes.add_shape(1, text_x, divider_top, Inches(0.35), Pt(1.5))
                divider.fill.solid()
                divider.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
                divider.line.fill.background()
                divider.shadow.inherit = False

                price_top = y + cell_h - Inches(0.55)
                desc_top = divider_top + Inches(0.17)
                # Same gap between the description box and the price as the
                # original fixed layout (0.4in) — desc_h is just whatever
                # space is left to reach that same fixed bottom edge, so a
                # blank model reproduces the exact original spacing, and a
                # present one shrinks the description's budget instead of
                # eating into the price's breathing room.
                desc_h = max(price_top - Inches(0.4) - desc_top, Inches(0.3))
                _add_text(slide, text_x, desc_top, text_w, desc_h,
                           product.description, 12, color=RGBColor(0x88, 0x88, 0x88), auto_fit=True)
                _add_text(slide, text_x, price_top, text_w, Inches(0.4),
                           f"{currency_symbol}{product.price:,.2f}", 18, bold=True)

                if getattr(product, "is_new", False):
                    _add_new_badge(slide, x, y, cell_w, accent_rgb)

    return prs
